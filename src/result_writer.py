import re
from typing import List

from pptx import Presentation
from pptx.util import Cm, Pt

from slide import Slide
from slide_image import SlideImage


class ResultWriter:
    def __init__(self):
        self.lines = 0

    def run(self, slides: List[Slide], out_file: str):
        prs = Presentation()
        prs.slide_height = Cm(14.288)
        prs.slide_width = Cm(25.4)
        tx_box = self.add_out_slide(prs)
        tf = tx_box.text_frame
        self.lines = 0
        for slide in slides:
            if self.lines > 30:
                tx_box = self.add_out_slide(prs)
                tf = tx_box.text_frame
            self.add_headline(tf, slide.slide_number, slide.title)
            self.add_images(tf, slide.images)
        prs.save(out_file)

    @staticmethod
    def add_out_slide(prs):
        blank_slide_layout = prs.slide_layouts[6]
        out_slide = prs.slides.add_slide(blank_slide_layout)
        left = Cm(2)  # 12
        top = Cm(1.75)
        height = Cm(11)
        width = Cm(9)
        return out_slide.shapes.add_textbox(left, top, width, height)

    def add_headline(self, tf, page: int, title: str):
        p = tf.add_paragraph()
        p.text = f"Page {page} - {title}"
        p.font.size = Pt(8)
        p.font.name = "Arial"
        p.font.bold = True
        self.lines += 1

    def add_images(self, tf, images: List[SlideImage]):
        lines = []
        for image in images:
            if image.status == SlideImage.STATUS_OK:
                for source in image.image_sources:
                    lines.append(source)
                for source in image.website_sources:
                    lines.append(source)
            else:
                lines.append(f"Error: [{image.status}]")
            lines.append("=")

        if lines[-1] == "=":
            lines.pop(-1)
        lines.append("\n")
        for line in lines:
            p = tf.add_paragraph()
            p.font.size = Pt(8)
            p.font.name = "Arial"

            if bool(re.match("^(\\[Error:\\s.+\\])|=|\n$", line)):
                p.text = line
            else:
                r = p.add_run()
                r.text = line
                hlink = r.hyperlink
                hlink.address = line
        self.lines += len(lines)
