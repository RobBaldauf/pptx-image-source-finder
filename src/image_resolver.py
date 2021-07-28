import io
import os
import re

from google.cloud import vision
from PIL import Image
from pptx import Presentation
from pptx.shapes.picture import Picture
from pptx.shapes.placeholder import PlaceholderPicture

from slide import Slide
from slide_image import SlideImage

SUPPORTED_EXTENSIONS = ["emf", "wmf", "jpg", "png", "gif", "jpeg", "tif"]


class ImageResolver:
    def __init__(self, key_filename: str):
        os.environ["GOOGLE_APPLICATION_CREDENTIALS"] = key_filename
        self.client = vision.ImageAnnotatorClient()

    def run(self, filename: str):
        pres = Presentation(filename)
        slides = []
        i = 1
        for slide in pres.slides:
            print(f"Processing slide {i}...")
            if slide.shapes.title:
                title = slide.shapes.title.text
                title = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f-\xff]", " ", title).replace("  ", " ")
            else:
                title = "???"
            slide_images = []
            for shape in slide.shapes:
                if isinstance(shape, (Picture, PlaceholderPicture)):
                    pic = shape._pic
                    img_rid = pic.xpath("./p:blipFill/a:blip/@r:embed")[0]
                    img_bytes = slide.part.related_parts[img_rid].image.blob
                    img_partname = slide.part.related_parts[img_rid].partname
                    slide_images.append(
                        self._get_sources(img_bytes, img_partname, slide.part.related_parts[img_rid].ext)
                    )
            if len(slide_images) > 0:
                slides.append(Slide(i, title, slide_images))
            i += 1
        return slides

    def _get_sources(self, img_bytes, img_partname: str, ext: str):
        # check extensions
        if ext.lower() not in SUPPORTED_EXTENSIONS:
            print(f"Error: Extension not supported: {ext}")
            return SlideImage(img_partname, [], [], SlideImage.STATUS_EXTENSION_NOT_SUPPORTED)
        img_bytes_local = img_bytes

        # handle emf/wmf files
        try:
            if ext in ["wmf", "emf"]:
                img_bytes_local = self._emf_to_png(img_bytes_local)
        except Exception as e:
            print(f"Error: Conversion failed...\n{e}")
            return SlideImage(img_partname, [], [], SlideImage.STATUS_CONVERSION_FAILED)

        # resize if necessary
        try:
            img_bytes_local = self._ensure_size(img_bytes_local)
        except Exception as e:
            print(f"Error: Resizing failed...\n{e}")
            return SlideImage(img_partname, [], [], SlideImage.STATUS_RESIZING_FAILED)

        # send request to API
        try:
            image_upload = vision.types.image_annotator_pb2.Image(content=img_bytes_local)
            response = self.client.annotate_image(
                {
                    "image": image_upload,
                    "features": [{"type": vision.enums.Feature.Type.WEB_DETECTION}],
                }
            )
        except Exception as e:
            print(f"Error: API request returned an error...\n{e}")
            return SlideImage(img_partname, [], [], SlideImage.STATUS_API_REQUEST_FAILED)

        # process response
        website_sources = []
        for x in response.web_detection.pages_with_matching_images:
            website_sources.append(x.url)
        image_sources = []
        for x in response.web_detection.full_matching_images:
            if "http" in x.url:
                image_sources.append(x.url)
        if len(image_sources) > 0 or len(website_sources) > 0:
            return SlideImage(img_partname, website_sources, image_sources, SlideImage.STATUS_OK)
        return SlideImage(img_partname, website_sources, image_sources, SlideImage.STATUS_NO_SOURCES_FOUND)

    @staticmethod
    def _emf_to_png(img_bytes):
        byte_io = io.BytesIO(img_bytes)
        img = Image.open(byte_io)
        buffered = io.BytesIO()
        img.save(buffered, format="PNG")
        return buffered.getvalue()

    @staticmethod
    def _ensure_size(img_bytes, max_width=1024):
        byte_io = io.BytesIO(img_bytes)
        img = Image.open(byte_io)

        if img.size[0] > max_width or img.size[1] > max_width:
            wpercent = max_width / float(img.size[0])
            hsize = int(float(img.size[1]) * float(wpercent))

            img = img.resize((max_width, hsize), Image.ANTIALIAS)

            buffered = io.BytesIO()
            if img.format is None:
                img.format = "PNG"
            img.save(buffered, format=img.format)
            return buffered.getvalue()
        return img_bytes
